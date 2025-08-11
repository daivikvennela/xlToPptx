import uuid
from flask import Flask, render_template, request, send_file, jsonify, Response
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from werkzeug.utils import secure_filename
import json
import logging
import numpy as np
from datetime import datetime
import shutil
from PIL import Image
import io
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
import tempfile
from pptx.dml.color import RGBColor
from docx import Document
import re
import csv
from docx.shared import RGBColor
from flask import request, jsonify
from lease_population.block_replacer import get_all_block_previews
from lease_population.block_replacer import replace_signature_and_notary_blocks
from docx import Document
from lease_population.block_replacer import generate_signature_block, generate_notary_block
from lease_population.block_replacer import getSigBlock, getNotaryBlock, generate_enhanced_combined_block

# --- Lease Population Module Integration ---
from lease_population import register_lease_population_routes


app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = None  # No size limit
app.config['MAX_CONTENT_PATH'] = None  # No path length limit
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0  # No caching for large files



def parse_kv_table_file(file_storage):
    """
    Parse a 2-column CSV or Excel file into a list of {'key': ..., 'value': ...} dicts.
    The first row is treated as the document name (not a key-value pair).
    """
    filename = file_storage.filename.lower()
    mapping = []
    document_name = None
    if filename.endswith('.csv'):
        file_storage.stream.seek(0)
        reader = csv.reader((line.decode('utf-8') for line in file_storage.stream), delimiter=',')
        rows = list(reader)
        if rows:
            document_name = rows[0][0].strip() if rows[0] and rows[0][0] else 'lease_population_filled'
            for row in rows[1:]:
                if len(row) >= 2:
                    mapping.append({'key': row[0].strip(), 'value': row[1].strip()})
    elif filename.endswith('.xlsx') or filename.endswith('.xls'):
        df = pd.read_excel(file_storage, header=None)
        if not df.empty:
            document_name = str(df.iloc[0,0]).strip() if not pd.isnull(df.iloc[0,0]) else 'lease_population_filled'
            for _, row in df.iloc[1:].iterrows():
                if len(row) >= 2:
                    mapping.append({'key': str(row[0]).strip(), 'value': str(row[1]).strip()})
    else:
        raise ValueError('Unsupported file type')
    # Remove header row if it looks like a header (after document name row)
    if mapping and mapping[0]['key'].lower() in ('key', 'placeholder') and mapping[0]['value'].lower() in ('value', 'replacement'):
        mapping = mapping[1:]
    return mapping, document_name

@app.route('/parse_kv_table', methods=['POST'])
def parse_kv_table():
    """
    Accept a 2-column CSV or Excel file and return a key-value mapping as JSON.
    The first row is treated as the document name.
    """
    if 'table_file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['table_file']
    try:
        mapping, document_name = parse_kv_table_file(file)
        # Validate
        keys = [pair['key'] for pair in mapping]
        if any(not k for k in keys) or len(set(keys)) != len(keys):
            return jsonify({'error': 'Keys must be non-empty and unique'}), 400
        return jsonify({'mapping': mapping, 'document_name': document_name})
    except Exception as e:
        return jsonify({'error': str(e)}), 400

# --- Track Changes Replacement Logic ---
def add_comment_to_run(run, comment_text):
    """
    Add a comment to a run. python-docx does not support comments natively as of 2024.
    This is a limitation. As a workaround, we can highlight the run and append the new value in brackets.
    """
    # Highlight the run (yellow)
    run.font.highlight_color = 7  # 7 = yellow in MS Word
    # Append the suggestion as [NEW: value]
    run.text += f" [NEW: {comment_text}]"
    # Note: True Word comments/revisions are not supported by python-docx as of 2024.
    # If/when support is added, this is where to insert a real comment.


def remove_acknowledgment_block(doc, keep_type):
    """
    Remove the acknowledgment block for the non-selected party type.
    - keep_type: 'Entity or Trust' or 'Individual'
    - The block to remove starts at its header and includes all following paragraphs up to the next block's header or end of document.
    - If both or neither block is found, insert a placeholder or error message.
    """
    # Define headers (case-insensitive, strip whitespace)
    header_entity = 'acknowledgment block for entity or trust'
    header_individual = 'acknowledgment block for individual'
    # Determine which to keep/remove
    if keep_type.lower() == 'entity or trust':
        header_remove = header_individual
        header_keep = header_entity
    else:
        header_remove = header_entity
        header_keep = header_individual
    # Find all paragraphs
    paragraphs = doc.paragraphs
    idx_remove = None
    idx_keep = None
    for i, p in enumerate(paragraphs):
        text = p.text.strip().lower()
        if text == header_remove:
            idx_remove = i
        if text == header_keep:
            idx_keep = i
    # If both or neither found, insert error/placeholder
    if idx_remove is None or idx_keep is None or idx_remove == idx_keep:
        doc.add_paragraph('[ERROR: Could not find both acknowledgment blocks for removal. Please check your template.]')
        return doc
    # Determine block to remove: from idx_remove up to (but not including) idx_keep or end
    start = idx_remove
    end = idx_keep if idx_keep > idx_remove else len(paragraphs)
    for i in range(end-1, start-1, -1):
        p = paragraphs[i]._element
        p.getparent().remove(p)
    return doc

def remove_entity_signature_block(doc):
    """
    Remove the '[Trust/Entity Name]' signature block and its following lines (By:, Name:, Title:) if present.
    """
    sig_header = '[trust/entity name]'
    sig_lines = ['by:', 'name:', 'title:']
    paragraphs = doc.paragraphs
    idx_sig = None
    for i, p in enumerate(paragraphs):
        if p.text.strip().lower() == sig_header:
            idx_sig = i
            break
    if idx_sig is not None:
        end = min(idx_sig + 4, len(paragraphs))
        for i in range(end-1, idx_sig-1, -1):
            p = paragraphs[i]._element
            p.getparent().remove(p)
    return doc

def remove_acknowledgment_blocks_enforced(doc, grantee_type):
    """
    Remove or retain sections based on grantee_type ('Individual' or 'Entity or Trust') using explicit start/end markers.
    For individuals:
      - Remove entity sections:
        1. [Trust/Entity Name] → My Commission Expires:___
        2. Acknowledgment Block for Entity or Trust → (Signature of Notary Public)
    For entities:
      - Remove individual sections:
        1. GRANTOR: → Name:
        2. Acknowledgment Block for Individual → (Signature of Notary Public)
    Only the relevant sections remain in the final document.
    If a marker is not found, return a clear error message.
    """
    grantee_type = grantee_type.strip().lower()
    paragraphs = doc.paragraphs
    def find_section_indices(start_marker, end_marker, section_label):
        indices = []
        start = None
        start_idx = None
        for i, p in enumerate(paragraphs):
            text = p.text.strip().lower()
            if start is None and start_marker.strip().lower() in text:
                start = i
                start_idx = i
            elif start is not None and end_marker.strip().lower() in text:
                indices.append((start, i))
                start = None
        if start is not None:
            context = '\n'.join(f'{j}: {paragraphs[j].text}' for j in range(max(0, start_idx-2), min(len(paragraphs), start_idx+5)))
            raise Exception(f"Could not find END marker '{end_marker}' for section '{section_label}'.\nContext:\n{context}")
        if not indices and start_marker:
            context = '\n'.join(f'{j}: {paragraphs[j].text}' for j in range(len(paragraphs)))
            raise Exception(f"Could not find START marker '{start_marker}' for section '{section_label}'.\nParagraphs:\n{context}")
        return indices
    to_remove = []
    if grantee_type == 'individual':
        # Remove entity sections
        entity1 = find_section_indices('[trust/entity name]', 'my commission expires:___', 'Entity Section 1')
        entity2 = find_section_indices('acknowledgment block for entity or trust', '(signature of notary public)', 'Entity Section 2')
        to_remove.extend(entity1)
        to_remove.extend(entity2)
    else:
        # Remove individual sections
        ind1 = find_section_indices('grantor:', 'name:', 'Individual Section 1')
        ind2 = find_section_indices('acknowledgment block for individual', '(signature of notary public)', 'Individual Section 2')
        to_remove.extend(ind1)
        to_remove.extend(ind2)
    for start, end in sorted(to_remove, reverse=True):
        for i in range(end, start-1, -1):
            p = paragraphs[i]._element
            p.getparent().remove(p)
    return doc

# --- Lease Population Placeholder Replacement: Robust Party Type Logic ---
# This endpoint now:
# 1. Always uses the current [Grantor Type] from the mapping to determine party type (case-insensitive, strict match).
# 2. Replaces all placeholders exactly as in the mapping, in all document sections (main, tables, headers, footers, footnotes), including [Signature Block] and [Notary Block].
# 3. Adds debug logging for party type and replaced placeholders.

# Lease population routes moved to lease_population.routes module

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/download/<filename>')
def download_file(filename):
    """Download generated presentation files"""
    try:
        file_path = os.path.join('uploads', filename)
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        print(f"Error downloading file: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/get_dynamic_block_preview', methods=['POST'])
def get_dynamic_block_preview():
    """Enhanced endpoint for dynamic signature/notary block preview with embedding logic"""
    try:
        data = request.get_json()
        
        # Extract form data
        owner_type = data.get('owner_type', 'individual')
        num_signatures = data.get('num_signatures', 1)
        state = data.get('state', '')
        county = data.get('county', '')
        name_of_individuals = data.get('name_of_individuals', '')
        type_of_authority = data.get('type_of_authority', '')
        instrument_for = data.get('instrument_for', '')
        grantor_name = data.get('grantor_name', '')
        trust_entity_name = data.get('trust_entity_name', '')
        name = data.get('name', '')
        title = data.get('title', '')
        include_signature = data.get('include_signature', True)
        include_notary = data.get('include_notary', True)
        embed_notary_in_signature = data.get('embed_notary_in_signature', True)
        
        # Generate enhanced combined block
        result = generate_enhanced_combined_block(
            owner_type=owner_type,
            grantor_name=grantor_name,
            trust_entity_name=trust_entity_name,
            name=name,
            title=title,
            state=state,
            county=county,
            name_of_individuals=name_of_individuals,
            type_of_authority=type_of_authority,
            instrument_for=instrument_for,
            num_signatures=num_signatures,
            include_signature=include_signature,
            include_notary=include_notary,
            embed_notary_in_signature=embed_notary_in_signature
        )
        
        return jsonify(result)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


    data = request.get_json()
    party_type = (data.get('party_type') or '').strip().lower()
    import os
    base = os.path.join('templates', 'blocks')
    if party_type == 'individual':
        sig_file = os.path.join(base, 'individual_signature.txt')
        notary_file = os.path.join(base, 'individual_notary.txt')
    else:
        sig_file = os.path.join(base, 'entity_signature.txt')
        notary_file = os.path.join(base, 'entity_notary.txt')
    try:
        with open(sig_file, 'r') as f:
            sig_block = f.read()
        with open(notary_file, 'r') as f:
            notary_block = f.read()
        return jsonify({'signature_block': sig_block, 'notary_block': notary_block})
    except Exception as e:
        return jsonify({'error': f'Could not load template: {str(e)}'}), 400

@app.route('/generate_signature_block', methods=['POST'])
def generate_signature_block():
    """
    Generate signature block using the generator function
    """
    try:
        data = request.get_json()
        owner_type = data.get('ownerType', '').strip()
        is_notary = data.get('isNotary', False)
        num_signatures = int(data.get('numSignatures', 1))
        
        print(f"[DEBUG] Generating signature block: owner_type='{owner_type}', is_notary={is_notary}, num_signatures={num_signatures}")
        
        # Validate inputs
        if not owner_type:
            return jsonify({'success': False, 'error': 'Owner type is required'}), 400
        
        if num_signatures < 1:
            return jsonify({'success': False, 'error': 'Number of signatures must be at least 1'}), 400
        
        # Import generator function
        from lease_population.block_replacer import generator
        
        # Call generator function
        result = generator(owner_type, is_notary, '', num_signatures)
        
        print(f"[DEBUG] Generated content length: {len(result) if result else 0}")
        
        return jsonify({'success': True, 'content': result})
        
    except Exception as e:
        print(f"[ERROR] Error generating signature block: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/get_signature_block', methods=['POST'])
def get_signature_block():
    data = request.get_json()
    owner_type = data.get('ownerType', '')
    num_signatures = data.get('numSignatures', 1)
    try:
        num_signatures = int(num_signatures)
    except Exception:
        num_signatures = 1
    content = getSigBlock(owner_type, num_signatures)
    return jsonify({'signature_block': content})

@app.route('/get_notary_block', methods=['POST'])
def get_notary_block():
    content = getNotaryBlock()
    return jsonify({'notary_block': content})

@app.route('/gen_exhibit_a', methods=['POST'])
def gen_exhibit_a():
    """
    Generate Exhibit A string from parcels, image, and description templates.
    """
    try:
        # Get parcels data
        parcels_data = request.form.get('parcels')
        if not parcels_data:
            return jsonify({'error': 'No parcels data provided'}), 400
        
        try:
            parcels = json.loads(parcels_data)
        except json.JSONDecodeError:
            return jsonify({'error': 'Invalid JSON format for parcels data'}), 400
        
        # Validate parcels data
        if not isinstance(parcels, list) or len(parcels) == 0:
            return jsonify({'error': 'Parcels data must be a non-empty list'}), 400
        
        print(f"[DEBUG] Processing {len(parcels)} parcels for Exhibit A generation")
        
        # Generate exhibit string using the build_exhibit_string function
        try:
            from lease_population.block_replacer import build_exhibit_string
            exhibit_string = build_exhibit_string(parcels)
            print(f"[DEBUG] Generated exhibit string, length: {len(exhibit_string)}")
            return jsonify({
                'exhibit_string': exhibit_string,
                'parcel_count': len(parcels)
            })
        except Exception as e:
            print(f"[ERROR] Failed to generate exhibit string: {str(e)}")
            import traceback
            traceback.print_exc()
            return jsonify({'error': f'Failed to generate exhibit string: {str(e)}'}), 500
        
    except Exception as e:
        import traceback
        error_traceback = traceback.format_exc()
        print(f"ERROR in gen_exhibit_a: {str(e)}")
        print(f"TRACEBACK: {error_traceback}")
        return jsonify({'error': f'Failed to generate exhibit string: {str(e)}', 'traceback': error_traceback}), 500

@app.route('/test_image_embedding', methods=['POST'])
def test_image_embedding():
    """
    Test route to verify image embedding functionality
    """
    try:
        if 'docx' not in request.files:
            return jsonify({'error': 'No DOCX file uploaded'}), 400
        
        docx_file = request.files['docx']
        if not docx_file.filename.endswith('.docx'):
            return jsonify({'error': 'Please upload a DOCX file'}), 400
        
        # Load the document
        doc = Document(docx_file)
        
        # Test with a sample image (1x1 pixel PNG)
        sample_image_data = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
        
        # Test the embedImage function
        from lease_population.block_replacer import embedImage
        success = embedImage(doc, sample_image_data, '[EXHIBIT_A_IMAGE_1]')
        
        if success:
            # Save the test document
            import tempfile
            import os
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                doc.save(tmp_file.name)
                tmp_file_path = tmp_file.name
            
            # Return success with file path
            return jsonify({
                'success': True,
                'message': 'Image embedding test successful',
                'file_path': tmp_file_path
            })
        else:
            return jsonify({'error': 'Image embedding test failed'}), 500
            
    except Exception as e:
        import traceback
        error_traceback = traceback.format_exc()
        print(f"ERROR in test_image_embedding: {str(e)}")
        print(f"TRACEBACK: {error_traceback}")
        return jsonify({'error': f'Test failed: {str(e)}', 'traceback': error_traceback}), 500

@app.route('/test_image_embedding_comprehensive', methods=['POST'])
def test_image_embedding_comprehensive():
    """
    Comprehensive test route for image embedding functionality.
    Tests various scenarios: valid PNG, invalid files, missing placeholders, etc.
    """
    try:
        print(f"[TEST] Starting comprehensive image embedding test")
        
        # Test 1: Valid PNG file
        if 'valid_png' in request.files:
            image_file = request.files['valid_png']
            print(f"[TEST] Testing valid PNG: {image_file.filename}")
            
            # Process the image
            image_data = image_file.read()
            if image_data.startswith(b'\x89PNG\r\n\x1a\n'):
                print(f"[TEST] ✓ Valid PNG header detected")
            else:
                print(f"[TEST] ✗ Invalid PNG header")
                return jsonify({'error': 'Invalid PNG file'}), 400
            
            # Convert to base64
            import base64
            img_b64 = base64.b64encode(image_data).decode('utf-8')
            print(f"[TEST] ✓ Base64 conversion successful: {len(img_b64)} chars")
            
            # Test image embedding function
            from lease_population.block_replacer import embedImage
            from docx import Document
            from io import BytesIO
            
            # Create test document
            doc = Document()
            doc.add_paragraph("Test document with [EXHIBIT_A_IMAGE_1] placeholder")
            
            # Test embedding
            success = embedImage(doc, img_b64, '[EXHIBIT_A_IMAGE_1]')
            if success:
                print(f"[TEST] ✓ Image embedding successful")
                
                # Save test document
                test_output = BytesIO()
                doc.save(test_output)
                test_output.seek(0)
                
                return jsonify({
                    'success': True,
                    'message': 'Image embedding test passed',
                    'image_size': len(image_data),
                    'base64_size': len(img_b64),
                    'embedding_success': True
                })
            else:
                print(f"[TEST] ✗ Image embedding failed")
                return jsonify({'error': 'Image embedding test failed'}), 500
        
        # Test 2: Invalid file type
        elif 'invalid_file' in request.files:
            print(f"[TEST] Testing invalid file type")
            return jsonify({'error': 'Invalid file type test - should be rejected'}), 400
        
        # Test 3: Missing placeholder
        elif 'missing_placeholder' in request.files:
            image_file = request.files['missing_placeholder']
            image_data = image_file.read()
            import base64
            img_b64 = base64.b64encode(image_data).decode('utf-8')
            
            # Create test document WITHOUT placeholder
            from docx import Document
            from lease_population.block_replacer import embedImage
            
            doc = Document()
            doc.add_paragraph("Test document WITHOUT placeholder")
            
            # Test embedding (should fail gracefully)
            success = embedImage(doc, img_b64, '[EXHIBIT_A_IMAGE_1]')
            if not success:
                print(f"[TEST] ✓ Correctly failed to embed image (no placeholder)")
                return jsonify({
                    'success': True,
                    'message': 'Missing placeholder test passed - correctly failed',
                    'embedding_success': False
                })
            else:
                print(f"[TEST] ✗ Should have failed but didn't")
                return jsonify({'error': 'Missing placeholder test failed'}), 500
        
        else:
            return jsonify({'error': 'No test file provided'}), 400
            
    except Exception as e:
        import traceback
        error_traceback = traceback.format_exc()
        print(f"[TEST ERROR] {str(e)}")
        print(f"[TEST TRACEBACK] {error_traceback}")
        return jsonify({'error': f'Test failed: {str(e)}', 'traceback': error_traceback}), 500

# Register lease population routes
register_lease_population_routes(app)

if __name__ == '__main__':
    import os
    port = int(os.environ.get('PORT', 5001))
    app.run(debug=True, port=port) 