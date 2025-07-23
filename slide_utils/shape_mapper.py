from pptx import Presentation
import json, pathlib, datetime

def build_map(pptx_path: str):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    m = {
        "source": pptx_path,
        "slide_index": 0,
        "generated": datetime.datetime.utcnow().isoformat(),
        "textboxes": []
    }
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        para = shp.text_frame.paragraphs[0] if shp.text_frame.paragraphs else None
        run = para.runs[0] if para and para.runs else None
        m["textboxes"].append({
            "shape_id": shp.shape_id,
            "name": getattr(shp, 'name', f"textbox_{shp.shape_id}"),
            "shape_type": "placeholder" if getattr(shp, 'is_placeholder', False) else "textbox",
            "bbox": [int(shp.left), int(shp.top), int(shp.width), int(shp.height)],
            "styles": {
                "font_name": run.font.name if run and run.font else None,
                "font_size": run.font.size.pt if run and run.font and run.font.size else None,
                "bold": bool(run.font.bold) if run and run.font and run.font.bold is not None else False,
                "italic": bool(run.font.italic) if run and run.font and run.font.italic is not None else False,
                "color_rgb": str(run.font.color.rgb) if run and run.font and run.font.color and hasattr(run.font.color, 'rgb') else None
            },
            "text_preview": shp.text[:80],
            "role": shp.placeholder_format.type if getattr(shp, 'is_placeholder', False) and hasattr(shp, 'placeholder_format') else None
        })
    map_path = pathlib.Path(pptx_path).with_suffix(".map.json")
    map_path.write_text(json.dumps(m, indent=2))
    return str(map_path) 