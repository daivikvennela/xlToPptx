from pptx import Presentation
import uuid, os

def inject_text(template_path: str, updates: list) -> str:
    prs = Presentation(template_path)
    slide = prs.slides[0]
    upd_map = {u["shape_id"]: u["new_text"] for u in updates}
    for shp in slide.shapes:
        if shp.shape_id in upd_map and shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            if shp.text_frame.paragraphs and shp.text_frame.paragraphs[0].runs:
                shp.text_frame.paragraphs[0].runs[0].text = upd_map[shp.shape_id]
    out_path = os.path.join("/tmp", f"{uuid.uuid4()}.pptx")
    prs.save(out_path)
    return out_path 